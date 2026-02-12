import nbformat
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import base64
import io
import os
from PIL import Image

# Paths
notebook_file = "alerts_analysis.ipynb"
ppt_file = "alerts_analysis_outputs.pptx"
img_folder = "temp_imgs"
os.makedirs(img_folder, exist_ok=True)

# Load notebook
with open(notebook_file, "r", encoding="utf-8") as f:
    nb = nbformat.read(f, as_version=4)

prs = Presentation()
img_count = 1

# === STEP 1: Add text and matplotlib outputs ===
for cell in nb.cells:
    if cell.cell_type == "code":
        for output in cell.get("outputs", []):
            slide_added = False  # Track if we add anything to this slide

            # 1️⃣ Matplotlib / PNG outputs
            if (
                output.output_type in ["display_data", "execute_result"]
                and "image/png" in output.data
            ):
                img_data = base64.b64decode(output.data["image/png"])
                img_path = os.path.join(img_folder, f"output_{img_count}.png")
                with open(img_path, "wb") as f:
                    f.write(img_data)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(
                    img_path, Inches(1), Inches(1), width=Inches(8)
                )
                img_count += 1
                slide_added = True

            # 2️⃣ Stream / printed text
            elif output.output_type == "stream" and output.name == "stdout":
                text = output.text.strip()
                if text:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    textbox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(6)
                    )
                    tf = textbox.text_frame
                    p = tf.add_paragraph()
                    p.text = text
                    p.font.size = Pt(18)
                    p.alignment = PP_ALIGN.LEFT
                    slide_added = True

            # 3️⃣ Plain text outputs (DataFrames, etc.)
            elif (
                output.output_type in ["execute_result", "display_data"]
                and "text/plain" in output.data
            ):
                text = output.data["text/plain"].strip()
                if text:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    textbox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(6)
                    )
                    tf = textbox.text_frame
                    p = tf.add_paragraph()
                    p.text = text
                    p.font.size = Pt(18)
                    p.alignment = PP_ALIGN.LEFT
                    slide_added = True

# === STEP 2: Add all saved Plotly images from temp_imgs folder ===
for file in sorted(os.listdir(img_folder)):
    if file.endswith(".png"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            os.path.join(img_folder, file), Inches(1), Inches(1), width=Inches(8)
        )

# Save presentation
prs.save(ppt_file)
print(f"Saved presentation: {ppt_file}")
